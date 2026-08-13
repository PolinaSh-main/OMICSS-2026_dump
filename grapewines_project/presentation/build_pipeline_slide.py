#!/usr/bin/env python3
"""
One 16:9 slide: how the data were treated, from raw VCF to candidate genes.

Every number on the slide is taken from the run that produced
results/ in this repository -- see PROVENANCE at the bottom of this
file for where each one comes from. Rebuild after changing it:

    python3 build_pipeline_slide.py

Output: pipeline_overview.pptx
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------
#  Palette
#
#  Lane colours come from the project palette in the README, so a lane
#  on this slide is the same colour as the figures it produces.
# ---------------------------------------------------------------------

INK = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x5A, 0x65, 0x70)
RULE = RGBColor(0x9E, 0xAD, 0xB6)

SLATE = RGBColor(0x45, 0x5A, 0x64)      # preprocessing, no lane of its own
BLUE = RGBColor(0x00, 0x69, 0x9A)       # K1 -- ancestry / FST
CRIMSON = RGBColor(0xAE, 0x00, 0x39)    # K5 -- ordination
GREEN = RGBColor(0x06, 0x40, 0x2B)      # K3 -- phylogeny

FONT = "Calibri"


def tint(colour: RGBColor, keep: float) -> RGBColor:
    """
    Blend towards white. keep=0.10 is a pale wash of the colour.
    """

    return RGBColor(*(
        round(channel * keep + 255 * (1 - keep))
        for channel in (colour[0], colour[1], colour[2])
    ))


# ---------------------------------------------------------------------
#  Geometry, in inches on a 13.333 x 7.5 slide
# ---------------------------------------------------------------------

SLIDE_W = 13.333
SLIDE_H = 7.5

COLUMN_W = 3.80

COLUMN_X = [0.45, 4.80, 9.15]           # left edge of each column
COLUMN_MID = [x + COLUMN_W / 2 for x in COLUMN_X]

BAND_Y = 1.28                           # preprocessing row
BAND_H = 1.06

BUS_Y = 2.62                            # horizontal distribution line

CHIP_Y = 2.90                           # lane header
CHIP_H = 0.30

BOX_Y = 3.32                            # first box of each lane
BOX_H = 1.10
BOX_GAP = 0.24

FOOTER_Y = 7.12
FOOTER_H = 0.26


def box_top(index: int) -> float:

    return BOX_Y + index * (BOX_H + BOX_GAP)


# ---------------------------------------------------------------------
#  Drawing helpers
# ---------------------------------------------------------------------


#
# Mini-markup for the two things plain runs cannot express: [[i:...]]
# for italic species names, [[sub:...]] for the ST in FST. Anything
# else is literal.
#

MARKUP = re.compile(r"\[\[(i|sub):(.*?)\]\]")


def fill_paragraph(paragraph, text, size, bold, colour):

    position = 0

    for match in MARKUP.finditer(text):

        if match.start() > position:
            _add_run(
                paragraph, text[position:match.start()],
                size, bold, colour
            )

        _add_run(
            paragraph, match.group(2), size, bold, colour,
            italic=(match.group(1) == "i"),
            subscript=(match.group(1) == "sub"),
        )

        position = match.end()

    if position < len(text):
        _add_run(paragraph, text[position:], size, bold, colour)


def _add_run(paragraph, text, size, bold, colour, italic=False, subscript=False):

    run = paragraph.add_run()

    run.text = text

    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic or None
    run.font.color.rgb = colour

    if subscript:
        # No API for this; baseline is a percentage in thousandths.
        run.font._rPr.set("baseline", "-25000")

    return run


def add_text(
    shapes,
    left, top, width, height,
    lines,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """
    lines: list of (text, size_pt, bold, colour, space_before_pt)
    """

    shape = shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    frame = shape.text_frame

    frame.word_wrap = True
    frame.vertical_anchor = anchor

    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for i, (text, size, bold, colour, space_before) in enumerate(lines):

        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()

        paragraph.alignment = align

        if space_before:
            paragraph.space_before = Pt(space_before)

        fill_paragraph(paragraph, text, size, bold, colour)

    return shape


def add_box(
    shapes, left, top, width, height, colour, title, detail,
    dashed=False,
):
    """
    A rounded card: pale wash of the lane colour, lane-coloured edge,
    dark title, grey detail line.
    """

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    shape.adjustments[0] = 0.10

    shape.fill.solid()
    shape.fill.fore_color.rgb = tint(colour, 0.07)

    shape.line.color.rgb = colour
    shape.line.width = Pt(1.5)

    if dashed:
        # Marks a box that is commentary rather than a processing step.
        shape.line._get_or_add_ln().append(
            shape.line._get_or_add_ln().makeelement(
                qn("a:prstDash"), {"val": "dash"}
            )
        )

    shape.shadow.inherit = False

    frame = shape.text_frame

    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)

    heading = frame.paragraphs[0]
    heading.alignment = PP_ALIGN.LEFT

    fill_paragraph(heading, title, 13.5, True, INK)

    body = frame.add_paragraph()
    body.alignment = PP_ALIGN.LEFT
    body.space_before = Pt(3)

    fill_paragraph(body, detail, 10.5, False, MUTED)

    return shape


def add_chip(shapes, left, top, width, height, colour, text):
    """
    Solid lane header. White on the lane colour -- every lane colour
    clears 4.5:1 against white.
    """

    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    shape.adjustments[0] = 0.28

    shape.fill.solid()
    shape.fill.fore_color.rgb = colour

    shape.line.fill.background()

    shape.shadow.inherit = False

    frame = shape.text_frame
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    frame.margin_top = 0
    frame.margin_bottom = 0

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(11.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return shape


def add_line(shapes, x1, y1, x2, y2, colour=RULE, width=1.75, arrow=True):

    connector = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )

    connector.line.color.rgb = colour
    connector.line.width = Pt(width)

    if arrow:
        # python-pptx has no arrowhead API; a:tailEnd is appended to the
        # line properties directly. It must come after the fill element,
        # which setting the colour above has already created.
        line_properties = connector.line._get_or_add_ln()

        tail = line_properties.makeelement(
            qn("a:tailEnd"),
            {"type": "triangle", "w": "med", "len": "med"}
        )

        line_properties.append(tail)

    return connector


# ---------------------------------------------------------------------
#  Content
# ---------------------------------------------------------------------

TITLE = "How the data were processed"

SUBTITLE = (
    "412 Caucasian grapevine accessions  ·  814,885 biallelic SNPs  ·  "
    "19 chromosomes  ·  8 metadata traits per accession"
)

PREPROCESSING = [
    (
        "1 · Raw variant calls",
        "cauca_grape.subset.vcf.gz, 412 accessions, plus the curated "
        "metadata table"
    ),
    (
        "2 · SNP filtering — vcftools",
        "PASS only · biallelic · call rate ≥ 0.6 · MAF ≥ 0.005"
    ),
    (
        "3 · PLINK conversion",
        "--make-bed --chr-set 19  →  814,885 SNPs × 412 samples"
    ),
]

LANES = [
    (
        BLUE,
        "ANCESTRY  →  DIFFERENTIATION",
        [
            (
                "ADMIXTURE, K = 2 … 10",
                "5-fold CV curve; K = 7 carried forward. 160 of 412 samples "
                "assigned at Q ≥ 0.75, 252 admixed"
            ),
            (
                "Pairwise F[[sub:ST]] — vcftools",
                "all 21 group pairs, per SNP; plain and weighted averages "
                "→ heatmap"
            ),
            (
                "Selection scan → genes",
                "50 kb windows (≥ 20 SNPs) · top 1% merged · genes from "
                "PN40024.v4.1"
            ),
        ],
    ),
    (
        CRIMSON,
        "ORDINATION",
        [
            (
                "PCA — plink --pca 20",
                "PC1 16.9%, PC2 9.7%, PC3 8.0% of the variance; scatter "
                "coloured by each trait"
            ),
            (
                "PC × metadata association",
                "one-way ANOVA per PC × trait; adjusted R², "
                "Benjamini–Hochberg FDR"
            ),
        ],
    ),
    (
        GREEN,
        "PHYLOGENY",
        [
            (
                "ML tree — SNPhylo",
                "MAF 0.10, LD 0.10 → 12,848 sites; rooted on ZZ01 "
                "([[i:V. rotundifolia]]); 100 bootstraps"
            ),
            (
                "NJ tree — sensitivity",
                "1 − IBS distances, LD r² < 0.2, midpoint rooted, "
                "100 bootstraps"
            ),
        ],
    ),
]

#
# Fills the space the two shorter lanes leave, and says why there are
# three of them rather than one. Drawn dashed so it does not read as a
# processing step.
#

COMMENTARY = (
    "Why three methods, not one",
    "ADMIXTURE assumes K ancestral populations · PCA assumes the structure "
    "is linear · the trees assume neither. Agreement between them is worth "
    "more than any one of them alone."
)


FOOTER = (
    "Nextflow on SLURM  ·  figures in Python (pandas / numpy / matplotlib)  ·  "
    "code and small results versioned in git"
)


def build(path: Path):

    presentation = Presentation()

    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    shapes = slide.shapes


    # ---- title -------------------------------------------------------

    add_text(
        shapes, 0.45, 0.24, 12.4, 0.52,
        [(TITLE, 30, True, INK, 0)]
    )

    add_text(
        shapes, 0.45, 0.80, 12.4, 0.34,
        [(SUBTITLE, 12.5, False, MUTED, 0)]
    )


    # ---- preprocessing row -------------------------------------------

    for i, (title, detail) in enumerate(PREPROCESSING):

        add_box(
            shapes,
            COLUMN_X[i], BAND_Y, COLUMN_W, BAND_H,
            SLATE, title, detail
        )

        if i:
            # Arrow into this box from the one on its left.
            gap_left = COLUMN_X[i - 1] + COLUMN_W
            add_line(
                shapes,
                gap_left + 0.08, BAND_Y + BAND_H / 2,
                COLUMN_X[i] - 0.06, BAND_Y + BAND_H / 2
            )


    # ---- distribution bus --------------------------------------------
    #
    # The three analyses all start from the same PLINK set. One trunk
    # down from step 3, one horizontal bus, one arrow down into each
    # lane -- so the slide says "parallel", not "sequential".

    add_line(
        shapes,
        COLUMN_MID[2], BAND_Y + BAND_H,
        COLUMN_MID[2], BUS_Y,
        arrow=False
    )

    add_line(
        shapes,
        COLUMN_MID[0], BUS_Y,
        COLUMN_MID[2], BUS_Y,
        arrow=False
    )

    for middle in COLUMN_MID:
        add_line(shapes, middle, BUS_Y, middle, CHIP_Y - 0.04)


    # ---- lanes -------------------------------------------------------

    for column, (colour, header, boxes) in enumerate(LANES):

        left = COLUMN_X[column]

        add_chip(
            shapes,
            left, CHIP_Y, COLUMN_W, CHIP_H,
            colour, header
        )

        for i, (title, detail) in enumerate(boxes):

            top = box_top(i)

            add_box(shapes, left, top, COLUMN_W, BOX_H, colour, title, detail)

            if i:
                add_line(
                    shapes,
                    left + COLUMN_W / 2, top - BOX_GAP + 0.02,
                    left + COLUMN_W / 2, top - 0.04
                )


    # ---- commentary, under the two shorter lanes ---------------------

    commentary_left = COLUMN_X[1]

    add_box(
        shapes,
        commentary_left, box_top(2),
        COLUMN_X[2] + COLUMN_W - commentary_left, BOX_H,
        SLATE, COMMENTARY[0], COMMENTARY[1],
        dashed=True
    )


    # ---- footer ------------------------------------------------------

    rule = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45), Inches(FOOTER_Y), Inches(12.44), Pt(1)
    )

    rule.fill.solid()
    rule.fill.fore_color.rgb = tint(SLATE, 0.30)
    rule.line.fill.background()
    rule.shadow.inherit = False

    add_text(
        shapes, 0.45, FOOTER_Y + 0.10, 12.44, FOOTER_H,
        [(FOOTER, 10.5, False, MUTED, 0)]
    )


    presentation.save(str(path))

    return path


if __name__ == "__main__":

    output = Path(__file__).resolve().parent / "pipeline_overview.pptx"

    build(output)

    print("Wrote", output)


# ---------------------------------------------------------------------
#  PROVENANCE
#
#  412 accessions, 19 chromosomes    reference/cauc_filtered.final.fam,
#                                    vcf_to_plink.sh --chr-set 19
#  814,885 SNPs                      results/pca/global/cauc_pca.log
#  filter thresholds                 00_preprocess/bin/filter_real_dataset.sh
#  K = 2..10, 5-fold CV              01_admixture/bin/run_admixture.sh
#  K = 7                             03_fst/nextflow.config -- the group's
#                                    choice; the CV error keeps falling to
#                                    K = 10, so do not call K = 7 a minimum
#  160 assigned / 252 admixed        results/fst/K7/sample_lists/
#                                    group_sizes_K7.tsv, at min_q 0.75
#  21 pairs                          7 choose 2
#  window 50 kb, >=20 SNPs, top 1%   04_manhattan_annotation/nextflow.config
#  PC1/PC2/PC3 variance              results/pca/global/explained_variance.tsv
#  12,848 sites, ZZ01, 100 boots     05_tree/README.md
#  LD r2 < 0.2, midpoint             05b_nj_tree/README.md
#  8 metadata traits                 reference/cauc_grape_metadata.csv has 10
#                                    columns: sample ID, variety name, 8 traits
# ---------------------------------------------------------------------
